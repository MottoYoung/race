import os
from pptx import Presentation
from PIL import Image
from moviepy.editor import *
import tempfile
import shutil
import nest_asyncio
import requests
import json
import base64
import hashlib
import time
from urllib.parse import urlencode, urlparse
import hmac
from datetime import datetime
from wsgiref.handlers import format_date_time
from time import mktime
import subprocess
import fitz  # PyMuPDF
nest_asyncio.apply()

class PPTtoVideo:
    def __init__(self, ppt_path, output_path, slide_duration=5, lang='zh-cn', use_generated_script=False, 
                 script_style="标准", target_audience="通用", total_duration="中等", subject="通用",
                 speed=50, vcn="x4_mingge", volume=50, pitch=50):
        """
        初始化PPT转视频类
        
        参数:
            ppt_path (str): PPT文件路径
            output_path (str): 输出视频路径
            slide_duration (int): 无语音笔记时每张幻灯片的默认时长（秒）
            lang (str): 文本转语音的语言
            use_generated_script (bool): 是否使用AI生成的演讲稿（True：使用AI生成，False：使用PPT备注）
            script_style (str): 演讲稿风格，可选"幽默"、"严谨"、"标准"、"激励"等
            target_audience (str): 目标受众，可选"小学生"、"初中生"、"高中生"、"大学生"、"专业人士"、"通用"
            total_duration (str): 演讲总时长，可选"简短"(3-5分钟)、"中等"(5-10分钟)、"详细"(10-15分钟)
            subject (str): 演讲学科，可选"语文"、"数学"、"英语"、"物理"、"化学"、"生物"、"历史"、"地理"、"政治"、"通用"等
            speed (int): 语速，取值范围[0-100]，默认50
            vcn (str): 发音人，默认"x4_mingge"（明哥-男声）
            volume (int): 音量，取值范围[0-100]，默认50
            pitch (int): 音调，取值范围[0-100]，默认50
        """
        self.ppt_path = ppt_path
        self.output_path = output_path
        self.slide_duration = slide_duration
        self.lang = lang
        self.use_generated_script = use_generated_script
        self.script_style = script_style
        self.target_audience = target_audience
        self.total_duration = total_duration
        self.subject = subject
        self.temp_dir = tempfile.mkdtemp()
        
        # 添加TTS语音合成参数
        self.speed = speed
        self.vcn = vcn
        self.volume = volume
        self.pitch = pitch
        
        # 根据total_duration参数设置演讲总时长范围
        self.duration_ranges = {
            "简短": "3-5分钟",
            "中等": "5-10分钟", 
            "详细": "10-15分钟"
        }
        
        # 讯飞API配置（用于语音合成）
        self.xf_host = "api-dx.xf-yun.com"
        self.xf_app_id = "b341104d"  # 替换为您的讯飞语音合成APPID
        self.xf_api_key = "569e475ac59475111c80757a163407b9"  # 替换为您的讯飞语音合成APIKey
        self.xf_api_secret = "YWI5ZDRhNTdmNGZjNmVjYjg3OWMwZGUz"  # 替换为您的讯飞语音合成APISecret
        
        # 星火大模型配置
        self.spark_url = "wss://spark-api.xf-yun.com/v4.0/chat"
        self.spark_app_id = "b341104d"  # 替换为您的星火大模型APPID
        self.spark_api_key = "569e475ac59475111c80757a163407b9"  # 替换为您的星火大模型APIKey
        self.spark_api_secret = "YWI5ZDRhNTdmNGZjNmVjYjg3OWMwZGUz"  # 替换为您的星火大模型APISecret
        self.domain = "4.0Ultra"
    
    def extract_slide_text(self, slide):
        """提取幻灯片中的文字内容"""
        text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                text.append(shape.text.strip())
        return "\n".join(filter(None, text))
        
    def get_script_from_notes(self, slide):
        """
        从PPT的备注中获取演讲稿
        """
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text.strip():
            return slide.notes_slide.notes_text_frame.text.strip()
        return ""

    def generate_complete_script(self, slides_text, total_slides):
        """
        使用星火大模型生成完整演讲稿，然后按幻灯片数量切割
        
        参数:
            slides_text (list): 所有幻灯片的文字内容列表
            total_slides (int): PPT总页数
        返回:
            list: 每页对应的演讲稿列表
        """
        print(f"\n正在为整个PPT（共{total_slides}页）生成完整演讲稿...")
        print(f"风格: {self.script_style}, 受众: {self.target_audience}, 总时长: {self.total_duration}, 学科: {self.subject}")
        
        from websocket import create_connection
        import json
        import ssl
        
        # 构建鉴权URL
        def create_url():
            # 使用与sparkAPI.py相同的鉴权逻辑
            host = urlparse(self.spark_url).netloc
            path = urlparse(self.spark_url).path
            
            # 生成RFC1123格式的时间戳
            now = datetime.now()
            date = format_date_time(mktime(now.timetuple()))
            
            # 拼接字符串
            signature_origin = f"host: {host}\ndate: {date}\nGET {path} HTTP/1.1"
            
            # 进行hmac-sha256进行加密
            signature_sha = hmac.new(self.spark_api_secret.encode('utf-8'),
                                   signature_origin.encode('utf-8'),
                                   digestmod=hashlib.sha256).digest()
            signature_sha_base64 = base64.b64encode(signature_sha).decode(encoding='utf-8')
            
            authorization_origin = f'api_key="{self.spark_api_key}", algorithm="hmac-sha256", headers="host date request-line", signature="{signature_sha_base64}"'
            authorization = base64.b64encode(authorization_origin.encode('utf-8')).decode(encoding='utf-8')
            
            # 将请求的鉴权参数组合为字典
            v = {
                "authorization": authorization,
                "date": date,
                "host": host
            }
            # 拼接鉴权参数
            url = self.spark_url + '?' + urlencode(v)
            return url

        # 构建所有幻灯片内容，更明确地标记页面边界
        all_slides_content = ""
        for i, text in enumerate(slides_text):
            # 添加明确的页面分隔符和页号
            all_slides_content += f"==================== 第{i+1}页 ====================\n"
            all_slides_content += text.strip() + "\n"
            all_slides_content += f"==================== 第{i+1}页结束 ====================\n\n"

        # 获取总时长范围
        total_duration_range = self.duration_ranges.get(self.total_duration, "5-10分钟")

        # 构建请求数据
        prompt_template = f"""你是一名优秀的{self.subject}学科教师，请根据以下{total_slides}页PPT的内容，生成一个完整的演讲稿。
我已经用明显的分隔符标记了每页PPT的开始和结束。

学科要求：{self.subject}学科的专业讲解
风格要求：{self.script_style}风格
目标受众：面向{self.target_audience}
时长要求：整个演讲稿的总时长控制在{total_duration_range}

具体要求：
1. 演讲稿必须分为{total_slides}个部分，每部分对应一页PPT
2. 各部分演讲稿的语言难度要适合{self.target_audience}理解
3. 各部分演讲稿的长度要合理分配，重要内容可以稍长，简单内容可以稍短
4. 第一部分应该有适当的开场白，不要过早总结整个演讲内容
5. 最后一部分应该有适当的总结和结束语
6. 各部分之间要有自然的过渡
7. 语言风格要符合"{self.script_style}"的特点
8. 使用{self.subject}学科的专业术语和教学方法

请严格按照如下格式返回演讲稿：

【第1页开始】
第1页的演讲内容...
【第1页结束】

【第2页开始】
第2页的演讲内容...
【第2页结束】

... 以此类推 ...

【第{total_slides}页开始】
第{total_slides}页的演讲内容...
【第{total_slides}页结束】

PPT内容：
{all_slides_content}
"""

        data = {
            "header": {
                "app_id": self.spark_app_id,
                "uid": "1234"
            },
            "parameter": {
                "chat": {
                    "domain": self.domain,
                    "temperature": 0.5,
                    "max_tokens": 8192  # 增加token上限，确保能生成完整内容
                }
            },
            "payload": {
                "message": {
                    "text": [
                        {"role": "user", "content": prompt_template}
                    ]
                }
            }
        }

        # 建立websocket连接
        print(f"正在连接星火API...")
        ws = create_connection(create_url(), sslopt={"cert_reqs": ssl.CERT_NONE})
        ws.send(json.dumps(data))
        print(f"请求已发送，等待响应...")

        # 获取返回结果
        complete_script = ""
        while True:
            result = ws.recv()
            result = json.loads(result)
            
            if result["header"]["code"] != 0:
                print(f"星火API调用失败：{result}")
                break
                
            choices = result["payload"]["choices"]
            content = choices["text"][0]["content"]
            complete_script += content
            
            # 判断是否是最后一条消息
            if choices["status"] == 2:
                print(f"完整演讲稿生成完成，长度: {len(complete_script)} 字符")
                break
            else:
                print(".", end="", flush=True)  # 显示进度
                
        ws.close()
        
        # 切割演讲稿
        script_parts = []
        for i in range(1, total_slides + 1):
            start_marker = f"【第{i}页开始】"
            end_marker = f"【第{i}页结束】"
            
            start_pos = complete_script.find(start_marker)
            if start_pos == -1:
                # 尝试其他可能的格式
                start_marker = f"【第 {i} 页开始】"
                start_pos = complete_script.find(start_marker)
                
                # 再尝试其他可能的格式
                if start_pos == -1:
                    alternative_markers = [
                        f"[第{i}页开始]", f"[第 {i} 页开始]",
                        f"<第{i}页开始>", f"<第 {i} 页开始>",
                        f"第{i}页:", f"第 {i} 页:"
                    ]
                    for marker in alternative_markers:
                        start_pos = complete_script.find(marker)
                        if start_pos != -1:
                            start_marker = marker
                            break
            
            end_pos = complete_script.find(end_marker)
            if end_pos == -1:
                # 尝试其他可能的格式
                end_marker = f"【第 {i} 页结束】"
                end_pos = complete_script.find(end_marker)
                
                # 再尝试其他可能的格式
                if end_pos == -1:
                    alternative_markers = [
                        f"[第{i}页结束]", f"[第 {i} 页结束]",
                        f"<第{i}页结束>", f"<第 {i} 页结束>"
                    ]
                    for marker in alternative_markers:
                        end_pos = complete_script.find(marker)
                        if end_pos != -1:
                            end_marker = marker
                            break
                            
                    # 如果仍然找不到结束标记，尝试查找下一页的开始标记作为当前页的结束
                    if end_pos == -1 and i < total_slides:
                        next_page_markers = [
                            f"【第{i+1}页开始】", f"【第 {i+1} 页开始】",
                            f"[第{i+1}页开始]", f"[第 {i+1} 页开始]",
                            f"<第{i+1}页开始>", f"<第 {i+1} 页开始>",
                            f"第{i+1}页:", f"第 {i+1} 页:"
                        ]
                        for marker in next_page_markers:
                            end_pos = complete_script.find(marker)
                            if end_pos != -1:
                                # 使用下一页的开始作为当前页的结束，但不包含标记本身
                                end_marker = ""
                                break
            
            if start_pos != -1 and end_pos != -1:
                # 提取该页的演讲稿（去掉标记）
                part_script = complete_script[start_pos + len(start_marker):end_pos].strip()
                script_parts.append(part_script)
                print(f"成功提取第{i}页演讲稿，长度: {len(part_script)} 字符")
            else:
                # 如果找不到标记，提供一个默认的简短描述
                print(f"警告：无法找到第{i}页的演讲稿标记，使用默认内容")
                if i == 1:
                    default_script = "让我们开始今天的演讲。"
                elif i == total_slides:
                    default_script = "感谢大家的聆听。"
                else:
                    default_script = f"这是第{i}页的内容。"
                script_parts.append(default_script)
        
        # 如果切割后的部分数量不足，补充空字符串
        while len(script_parts) < total_slides:
            script_parts.append("")
        
        return script_parts

    def extract_slides_and_notes(self):
        """提取PPT中的幻灯片和笔记"""
        prs = Presentation(self.ppt_path)
        slides_data = []
        
        # 获取总页数
        total_slides = len(prs.slides)
        
        # 首先将整个PPT转换为PDF
        pdf_path = os.path.join(self.temp_dir, "slides.pdf")
        self._convert_ppt_to_pdf(self.ppt_path, pdf_path)
        
        # 将PDF转换为图片
        images = self._convert_pdf_to_images(pdf_path)
        
        # 如果使用AI生成演讲稿，先提取所有幻灯片的文字内容
        if self.use_generated_script:
            all_slides_text = []
            for slide in prs.slides:
                slide_text = self.extract_slide_text(slide)
                all_slides_text.append(slide_text)
            
            # 生成完整演讲稿并切割
            script_parts = self.generate_complete_script(all_slides_text, total_slides)
        
        for i, slide in enumerate(prs.slides):
            # 保存幻灯片为图片
            if i < len(images):
                slide_path = os.path.join(self.temp_dir, f"slide_{i}.png")
                images[i].save(slide_path, "PNG")
            else:
                # 如果PDF转换失败，使用备用方法
                slide_path = os.path.join(self.temp_dir, f"slide_{i}.png")
                self._save_slide_as_image_fallback(slide, slide_path)
            
            # 根据参数决定使用备注还是生成演讲稿
            notes = ""
            if self.use_generated_script:
                # 使用预先生成的完整演讲稿中对应的部分
                if i < len(script_parts):
                    notes = script_parts[i]
                # 如果生成的演讲稿为空但有备注，则使用备注
                if not notes.strip() and slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                    notes = slide.notes_slide.notes_text_frame.text
            else:
                # 优先使用PPT备注
                if slide.has_notes_slide and slide.notes_slide.notes_text_frame.text:
                    notes = slide.notes_slide.notes_text_frame.text
            
            slides_data.append({
                'image_path': slide_path,
                'notes': notes
            })
        
        return slides_data
    
    def _convert_ppt_to_pdf(self, ppt_path, pdf_path):
        """使用LibreOffice将PPT转换为PDF"""
        try:
            # 尝试可能的LibreOffice路径
            libreoffice_paths = [
                "soffice",  # 如果在PATH中
                "C:\\Program Files\\LibreOffice\\program\\soffice.exe",  # 标准安装路径
                "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",  # 32位系统路径
            ]
            
            # 如果有Microsoft Office，也可以尝试使用PowerPoint转换
            ms_office_paths = [
                "C:\\Program Files\\Microsoft Office\\root\\Office16\\POWERPNT.EXE",
                "C:\\Program Files (x86)\\Microsoft Office\\root\\Office16\\POWERPNT.EXE",
                "C:\\Program Files\\Microsoft Office\\Office16\\POWERPNT.EXE",
                "C:\\Program Files (x86)\\Microsoft Office\\Office16\\POWERPNT.EXE",
            ]
            
            # 首先尝试使用LibreOffice
            success = False
            error_msg = ""
            
            for libreoffice_cmd_path in libreoffice_paths:
                if shutil.which(libreoffice_cmd_path) or os.path.exists(libreoffice_cmd_path):
                    libreoffice_cmd = [
                        libreoffice_cmd_path,
                        "--headless",
                        "--convert-to", "pdf",
                        "--outdir", os.path.dirname(pdf_path),
                        ppt_path
                    ]
                    
                    try:
                        process = subprocess.Popen(
                            libreoffice_cmd, 
                            stdout=subprocess.PIPE, 
                            stderr=subprocess.PIPE
                        )
                        stdout, stderr = process.communicate()
                        
                        if process.returncode == 0:
                            # 由于LibreOffice会使用原文件名，可能需要重命名
                            base_name = os.path.basename(ppt_path)
                            name_without_ext = os.path.splitext(base_name)[0]
                            original_output = os.path.join(os.path.dirname(pdf_path), f"{name_without_ext}.pdf")
                            
                            if original_output != pdf_path and os.path.exists(original_output):
                                shutil.move(original_output, pdf_path)
                            
                            success = os.path.exists(pdf_path)
                            if success:
                                print(f"使用 {libreoffice_cmd_path} 成功转换PPT到PDF")
                                return True
                        else:
                            error_msg += f"使用 {libreoffice_cmd_path} 失败: {stderr.decode('utf-8', errors='ignore')}\n"
                    except Exception as e:
                        error_msg += f"使用 {libreoffice_cmd_path} 出错: {str(e)}\n"
            
            # 如果LibreOffice失败，尝试使用comtypes和PowerPoint (仅限Windows)
            if not success and os.name == 'nt':
                try:
                    import comtypes.client
                    
                    # 创建PowerPoint应用实例
                    powerpoint = comtypes.client.CreateObject("Powerpoint.Application")
                    powerpoint.Visible = True
                    
                    # 打开PPT文件
                    presentation = powerpoint.Presentations.Open(os.path.abspath(ppt_path))
                    
                    # 导出为PDF
                    presentation.ExportAsFixedFormat(os.path.abspath(pdf_path), 
                                                     32,  # ppFixedFormatTypePDF
                                                     PrintRange=None,
                                                     OutputType=0,  # ppPrintOutputSlides
                                                     PrintHiddenSlides=False)
                    
                    # 关闭文件和应用
                    presentation.Close()
                    powerpoint.Quit()
                    
                    success = os.path.exists(pdf_path)
                    if success:
                        print("成功使用PowerPoint转换PPT到PDF")
                        return True
                except Exception as e:
                    error_msg += f"使用PowerPoint COM转换失败: {str(e)}\n"
            
            if not success:
                print(f"PPT转PDF失败，所有方法都失败: \n{error_msg}")
                return False
                
        except Exception as e:
            print(f"PPT转PDF失败: {str(e)}")
            return False
    
    def _convert_pdf_to_images(self, pdf_path, dpi=200):
        """使用PyMuPDF将PDF转换为图片列表"""
        try:
            if not os.path.exists(pdf_path):
                print(f"PDF文件不存在: {pdf_path}")
                return []
                
            # 计算缩放因子（dpi/72，因为PDF使用72dpi作为基础单位）
            zoom = dpi / 72
            
            # 打开PDF文件
            pdf_document = fitz.open(pdf_path)
            images = []
            
            # 遍历每一页
            for page_num in range(len(pdf_document)):
                page = pdf_document.load_page(page_num)
                
                # 渲染页面为像素图
                pix = page.get_pixmap(matrix=fitz.Matrix(zoom, zoom))
                
                # 转换为PIL图像
                img = Image.frombytes("RGB", [pix.width, pix.height], pix.samples)
                
                # 确保宽度和高度是偶数
                width, height = img.size
                if width % 2 != 0 or height % 2 != 0:
                    # 计算新的偶数尺寸
                    new_width = width if width % 2 == 0 else width + 1
                    new_height = height if height % 2 == 0 else height + 1
                    # 调整图像大小
                    img = img.resize((new_width, new_height), Image.LANCZOS)
                    print(f"调整图像尺寸从 {width}x{height} 到 {new_width}x{new_height}")
                
                images.append(img)
            
            pdf_document.close()
            return images
        except Exception as e:
            print(f"PDF转图片失败: {str(e)}")
            return []
    
    def _save_slide_as_image_fallback(self, slide, output_path, width=1280, height=720):
        """备用方法：当PDF转换失败时创建占位图像"""
        # 确保宽度和高度是偶数
        width = width if width % 2 == 0 else width + 1
        height = height if height % 2 == 0 else height + 1
        
        img = Image.new('RGB', (width, height), color='white')
        img.save(output_path)
        print(f"使用备用方法创建了占位图像: {output_path}")
    
    def _assemble_auth_url(self, path):
        """生成讯飞API鉴权URL"""
        params = self._assemble_auth_params(path)
        request_url = "http://" + self.xf_host + path
        auth_url = request_url + "?" + urlencode(params)
        return auth_url
    
    def _assemble_auth_params(self, path):
        """生成讯飞API鉴权参数"""
        format_date = format_date_time(mktime(datetime.now().timetuple()))
        signature_origin = "host: " + self.xf_host + "\n"
        signature_origin += "date: " + format_date + "\n"
        signature_origin += "POST " + path + " HTTP/1.1"
        
        signature_sha = hmac.new(self.xf_api_secret.encode('utf-8'), 
                                 signature_origin.encode('utf-8'),
                                 digestmod=hashlib.sha256).digest()
        signature_sha = base64.b64encode(signature_sha).decode(encoding='utf-8')
        
        authorization_origin = 'api_key="%s", algorithm="%s", headers="%s", signature="%s"' % (
            self.xf_api_key, "hmac-sha256", "host date request-line", signature_sha)
        authorization = base64.b64encode(authorization_origin.encode('utf-8')).decode(encoding='utf-8')
        
        params = {
            "host": self.xf_host,
            "date": format_date,
            "authorization": authorization
        }
        return params
    
    def create_audio_from_notes(self, notes, output_path):
        """将笔记转换为语音"""
        if notes.strip():
            try:
                # 创建任务
                task_id = self._create_tts_task(notes)
                if not task_id:
                    print("讯飞TTS任务创建失败")
                    return False
                
                # 查询任务结果
                audio_url = self._query_tts_task(task_id)
                if not audio_url:
                    print("讯飞TTS任务查询失败")
                    return False
                
                # 下载音频文件
                response = requests.get(audio_url)
                if response.status_code == 200:
                    with open(output_path, 'wb') as f:
                        f.write(response.content)
                    return os.path.exists(output_path)
                else:
                    print(f"音频下载失败，状态码: {response.status_code}")
                    return False
            except Exception as e:
                print(f"语音生成失败: {str(e)}")
                # 生成空音频文件防止后续崩溃
                open(output_path, 'w').close()
                return False
        return False
    
    def _create_tts_task(self, text):
        """创建讯飞TTS任务"""
        create_path = "/v1/private/dts_create"
        auth_url = self._assemble_auth_url(create_path)
        
        # 文本内容Base64编码
        encode_str = base64.b64encode(text.encode("UTF8"))
        txt = encode_str.decode()
        
        # 请求头和请求内容
        headers = {'Content-Type': 'application/json'}
        data = {
            "header": {
                "app_id": self.xf_app_id
            },
            "parameter": {
                "dts": {
                    "vcn": self.vcn,  # 使用实例变量
                    "language": "zh",
                    "speed": self.speed,  # 使用实例变量
                    "volume": self.volume,  # 使用实例变量
                    "pitch": self.pitch,  # 使用实例变量
                    "rhy": 0,
                    "audio": {
                        "encoding": "lame",  # MP3格式
                        "sample_rate": 16000,
                        "channels": 1,
                        "bit_depth": 16
                    },
                    "pybuf": {
                        "encoding": "utf8",
                        "compress": "raw",
                        "format": "plain"
                    }
                }
            },
            "payload": {
                "text": {
                    "encoding": "utf8",
                    "compress": "raw",
                    "format": "plain",
                    "text": txt
                }
            }
        }
        
        try:
            response = requests.post(url=auth_url, headers=headers, data=json.dumps(data))
            result = json.loads(response.text)
            
            if result.get('header', {}).get('code') == 0:
                task_id = result.get('header', {}).get('task_id')
                print(f"讯飞TTS任务创建成功，任务ID: {task_id}")
                return task_id
            else:
                code = result.get('header', {}).get('code')
                msg = result.get('header', {}).get('message')
                print(f"讯飞TTS任务创建失败，错误码: {code}, 错误信息: {msg}")
                return None
        except Exception as e:
            print(f"讯飞TTS任务创建异常: {str(e)}")
            return None
    
    def _query_tts_task(self, task_id):
        """查询讯飞TTS任务结果"""
        query_path = "/v1/private/dts_query"
        auth_url = self._assemble_auth_url(query_path)
        
        headers = {'Content-Type': 'application/json'}
        data = {
            "header": {
                "app_id": self.xf_app_id,
                "task_id": task_id
            }
        }
        
        # 循环查询，最多查询20次
        for i in range(20):  # 将查询次数从10次增加到20次
            try:
                time.sleep(2)  # 等待时间从1秒增加到2秒
                response = requests.post(url=auth_url, headers=headers, data=json.dumps(data))
                result = json.loads(response.text)
                
                code = result.get('header', {}).get('code')
                if code == 0:
                    task_status = result.get('header', {}).get('task_status')
                    if task_status == '5':  # 任务处理成功
                        audio_base64 = result.get('payload', {}).get('audio', {}).get('audio')
                        if audio_base64:
                            audio_url = base64.b64decode(audio_base64).decode()
                            print(f"讯飞TTS任务查询成功，音频URL: {audio_url}")
                            return audio_url
                    elif task_status in ['1', '3']:  # 任务创建成功或处理中
                        print(f"讯飞TTS任务处理中，第{i+1}次查询")
                        continue
                    else:  # 任务派发失败或处理失败
                        print(f"讯飞TTS任务失败，状态码: {task_status}")
                        return None
                else:
                    print(f"讯飞TTS任务查询失败，错误码: {code}")
                    return None
            except Exception as e:
                print(f"讯飞TTS任务查询异常: {str(e)}")
                return None
        
        print("讯飞TTS任务查询超时")
        return None
    
    def create_video(self):
        """创建视频"""
        slides_data = self.extract_slides_and_notes()
        clips = []
        
        for i, slide in enumerate(slides_data):
            # 创建语音
            audio_path = os.path.join(self.temp_dir, f"audio_{i}.mp3")
            has_audio = self.create_audio_from_notes(slide['notes'], audio_path)
            
            # 创建图像剪辑
            img_clip = ImageClip(slide['image_path'])
            
            # 确保图像尺寸是偶数
            img_size = img_clip.size
            if img_size[0] % 2 != 0 or img_size[1] % 2 != 0:
                new_size = (
                    img_size[0] if img_size[0] % 2 == 0 else img_size[0] + 1,
                    img_size[1] if img_size[1] % 2 == 0 else img_size[1] + 1
                )
                img_clip = img_clip.resize(newsize=new_size)
                print(f"调整剪辑尺寸从 {img_size} 到 {new_size}")
            
            # 音频处理
            if has_audio:
                try:
                    audio_clip = AudioFileClip(audio_path)
                    duration = audio_clip.duration
                    img_clip = img_clip.set_duration(duration)
                    img_clip = img_clip.set_audio(audio_clip)
                except Exception as e:
                    print(f"音频处理失败: {str(e)}，使用默认时长")
                    img_clip = img_clip.set_duration(self.slide_duration)
            else:
                # 如果没有语音，使用默认持续时间
                img_clip = img_clip.set_duration(self.slide_duration)
            
            clips.append(img_clip)
        
        # 合并所有剪辑
        final_clip = concatenate_videoclips(clips)
        
        # 导出视频 - 使用更兼容的编码参数
        final_clip.write_videofile(
            self.output_path, 
            fps=24, 
            codec='libx264', 
            audio_codec='aac',
            preset='ultrafast',  # 加快编码速度
            ffmpeg_params=["-pix_fmt", "yuv420p"]  # 增加兼容性
        )
        
        # 清理临时文件
        self.cleanup()
        
    def cleanup(self):
        """清理临时文件"""
        shutil.rmtree(self.temp_dir)

def convert_ppt_to_video(ppt_path, output_path, slide_duration=5, lang='zh-cn', use_generated_script=False,
                        script_style="标准", target_audience="通用", total_duration="中等", subject="通用",
                        speed=50, vcn="x4_mingge", volume=50, pitch=50):
    """
    将PPT转换为视频的便捷函数
    
    参数:
        ppt_path (str): PPT文件路径
        output_path (str): 输出视频路径
        slide_duration (int): 无语音笔记时每张幻灯片的默认时长（秒）
        lang (str): 文本转语音的语言（默认为中文）
        use_generated_script (bool): 是否使用AI生成的演讲稿（True：使用AI生成，False：使用PPT备注）
        script_style (str): 演讲稿风格，可选"幽默"、"严谨"、"标准"、"激励"等
        target_audience (str): 目标受众，可选"小学生"、"初中生"、"高中生"、"大学生"、"专业人士"、"通用"
        total_duration (str): 演讲总时长，可选"简短"(3-5分钟)、"中等"(5-10分钟)、"详细"(10-15分钟)
        subject (str): 演讲学科，可选"语文"、"数学"、"英语"、"物理"、"化学"、"生物"、"历史"、"地理"、"政治"、"通用"等
        speed (int): 语速，取值范围[0-100]，默认50
        vcn (str): 发音人，默认"x4_mingge"（明哥-男声）
        volume (int): 音量，取值范围[0-100]，默认50
        pitch (int): 音调，取值范围[0-100]，默认50
    """
    converter = PPTtoVideo(ppt_path, output_path, slide_duration, lang, use_generated_script,
                          script_style, target_audience, total_duration, subject,
                          speed, vcn, volume, pitch)
    converter.create_video()
    print(f"视频已生成到: {output_path}")

# 使用示例
if __name__ == "__main__":
    # 直接在此处指定路径和参数
    ppt_path = "D:/CODE/demo_ex/Now2/input.pptx"
    output_path = "D:/CODE/demo_ex/test/output.mp4"
    duration = 5
    
    # 示例3：生成适合小学生的幽默风格数学课演讲
    convert_ppt_to_video(
        ppt_path=ppt_path,
        output_path="output_math_for_czs.mp4",
        slide_duration=duration,
        use_generated_script=True,
        script_style="严谨",
        target_audience="初中生",
        total_duration="简短",
        subject="历史",
        speed=50,  # 默认语速
        vcn="x4_mingge",  # 默认发音人
        volume=50,  # 默认音量
        pitch=50  # 默认音调
    )
    print("历史")
    
    # 示例4：生成适合大学生的严谨物理学演讲
    # convert_ppt_to_video(
    #     ppt_path=ppt_path,
    #     output_path="output_physics_university.mp4",
    #     slide_duration=duration,
    #     use_generated_script=True,
    #     script_style="严谨",
    #     target_audience="大学生",
    #     total_duration="详细",
    #     subject="物理",
    #     speed=50,
    #     vcn="x4_lingbosong",
    #     volume=50,
    #     pitch=50
    # )
    # print("已生成适合大学生的严谨物理学演讲视频")
